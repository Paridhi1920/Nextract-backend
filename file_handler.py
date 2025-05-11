import docx
from pptx import Presentation
import pdfplumber

def extract_text_from_pdf(file_path):
    try:
        with pdfplumber.open(file_path) as pdf:
            text = ""
            for page in pdf.pages:
                try:
                    text += page.extract_text() + "\n"
                except Exception as e:
                    print(f"Error reading page: {e}")
            return text if text.strip() != "" else "No text extracted."
    except Exception as e:
        print(f"PDF extraction failed: {e}")
        return "PDF extraction failed."

def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text.strip()
